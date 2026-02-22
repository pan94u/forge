#!/usr/bin/env python3
"""Generate Forge Platform leadership presentation from baseline-v1.4."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)    # Primary blue
ACCENT_GREEN = RGBColor(0x00, 0xB4, 0x8A)   # Success green
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)  # Warning orange
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)     # Light background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
SECTION_BG = RGBColor(0xE8, 0xF4, 0xFD)     # Light blue section bg
RED = RGBColor(0xE0, 0x40, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        # Set transparency via XML
        import lxml.etree as etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_elem.set('val', str(int(alpha * 1000)))
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_para(tf, text, size=16, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4), font_name='Microsoft YaHei'):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_bullet(tf, text, level=0, size=16, bold=False, color=TEXT_DARK):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Microsoft YaHei'
    return p


def add_section_header(slide):
    """Add a thin accent bar at the top."""
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)


def make_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)

    # Accent bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.12), ACCENT_BLUE)

    # Title
    tb = add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5))
    set_text(tb.text_frame, 'Forge 智能交付平台', size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame, '面向 AI 时代的全生命周期交付系统', size=24, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER, space_before=Pt(12))

    # Subtitle box
    box = add_shape(slide, Inches(3.5), Inches(4.2), Inches(6.3), Inches(0.7), ACCENT_BLUE)
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = box.text_frame.paragraphs[0].add_run()
    run.text = '规划基线 v1.4  |  2026-02-19'
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = 'Microsoft YaHei'

    # Stats bar
    stats = [
        ('320+', '文件'),
        ('45,000+', '行代码'),
        ('130+', '测试用例'),
        ('9', 'MCP 工具'),
        ('4', '容器部署'),
    ]
    start_x = Inches(1.5)
    for i, (num, label) in enumerate(stats):
        x = start_x + Inches(i * 2.1)
        tb = add_textbox(slide, x, Inches(5.4), Inches(1.8), Inches(1.2))
        set_text(tb.text_frame, num, size=32, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
        add_para(tb.text_frame, label, size=14, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)


def make_goals_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '平台目标与定位', size=32, bold=True, color=DARK_BG)

    goals = [
        ('交付提效', 'SuperAgent + Skill 体系，5-7 人团队 → 1-2 人 + AI 协同'),
        ('知识沉淀', '隐性知识编码为可复用的 Skill 和结构化数据'),
        ('持续进化', '交付环 + 进化环双环机制，平台在使用中自动变好'),
        ('全员可用', 'CLI/IDE（开发者）+ Web IDE（产品团队）统一入口'),
        ('跨栈迁移', 'AI 知识抢救 + 跨语言翻译，10-13x 人天提效'),
    ]

    colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE, RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xE0, 0x40, 0x40)]

    for i, (title, desc) in enumerate(goals):
        y = Inches(1.5) + Inches(i * 1.1)
        # Colored left bar
        add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(0.85), colors[i])
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.15), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = circle.text_frame.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = 'Microsoft YaHei'
        circle.text_frame.word_wrap = False

        # Title + desc
        tb = add_textbox(slide, Inches(1.9), y, Inches(10), Inches(0.85))
        set_text(tb.text_frame, title, size=20, bold=True, color=TEXT_DARK)
        add_para(tb.text_frame, desc, size=16, color=TEXT_GRAY, space_before=Pt(2))

    # Right side: positioning box
    box = add_shape(slide, Inches(9.5), Inches(1.3), Inches(3.5), Inches(5.5), SECTION_BG)
    tb = add_textbox(slide, Inches(9.7), Inches(1.5), Inches(3.1), Inches(5))
    set_text(tb.text_frame, '核心判断', size=18, bold=True, color=ACCENT_BLUE)
    bullets = [
        'Skill 是核心竞争力',
        'Code is the universal interface',
        '双环 > 单环',
        'Runtime 无关是长期战略',
        '质量底线 > 开发速度',
    ]
    for b in bullets:
        add_bullet(tb.text_frame, '• ' + b, size=14, color=TEXT_DARK)


def make_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '双环架构 + SuperAgent', size=32, bold=True, color=DARK_BG)

    # Loop 1 box
    box1 = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.5), RGBColor(0xE3, 0xF2, 0xFD))
    tb1 = add_textbox(slide, Inches(1), Inches(1.6), Inches(5.4), Inches(2.3))
    set_text(tb1.text_frame, '环 1：交付环 (Delivery Loop)', size=20, bold=True, color=ACCENT_BLUE)
    add_para(tb1.text_frame, '', size=8)
    add_bullet(tb1.text_frame, '人类指定 WHAT（声明式意图）', size=15)
    add_bullet(tb1.text_frame, 'SuperAgent 解决 HOW（命令式执行）', size=15)
    add_bullet(tb1.text_frame, 'HITL 人在关键节点审批', size=15)
    add_para(tb1.text_frame, '', size=6)
    add_para(tb1.text_frame, '规划 → 设计 → 开发 → 测试 → 运维 → 运营', size=14, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Loop 2 box
    box2 = add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.8), Inches(2.5), RGBColor(0xE8, 0xF5, 0xE9))
    tb2 = add_textbox(slide, Inches(1), Inches(4.4), Inches(5.4), Inches(2.3))
    set_text(tb2.text_frame, '环 2：进化环 (Learning Loop)', size=20, bold=True, color=ACCENT_GREEN)
    add_para(tb2.text_frame, '', size=8)
    add_bullet(tb2.text_frame, '每次交付 → 生产数据 → 资产提取', size=15)
    add_bullet(tb2.text_frame, '私域知识沉淀 → 反哺 SuperAgent', size=15)
    add_bullet(tb2.text_frame, '知识库自然增长 / Skill 自动优化', size=15)
    add_para(tb2.text_frame, '', size=6)
    add_para(tb2.text_frame, '⏳ Phase 3 实现', size=14, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

    # SuperAgent box (right)
    sa_box = add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xF3, 0xE5, 0xF5))
    tb3 = add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb3.text_frame, 'SuperAgent（唯一实例）', size=20, bold=True, color=RGBColor(0x8B, 0x5C, 0xF6))
    add_para(tb3.text_frame, '', size=8)
    add_para(tb3.text_frame, 'Skill Profile Router', size=16, bold=True, color=TEXT_DARK)
    add_para(tb3.text_frame, '根据交付阶段 / 用户指定 / AI 自主判断切换 Profile', size=13, color=TEXT_GRAY, space_before=Pt(2))
    add_para(tb3.text_frame, '', size=6)

    profiles = ['规划 Profile', '设计 Profile', '开发 Profile', '测试 Profile', '运维 Profile']
    for p in profiles:
        add_bullet(tb3.text_frame, '✅ ' + p, size=14, color=TEXT_DARK)

    add_para(tb3.text_frame, '', size=8)
    add_para(tb3.text_frame, 'OODA 内循环', size=16, bold=True, color=TEXT_DARK)
    add_para(tb3.text_frame, 'Observe → Orient → Decide → Act', size=14, color=ACCENT_BLUE)
    add_para(tb3.text_frame, '底线 ✅ → HITL 审批  |  底线 ❌ → 自动修复', size=13, color=TEXT_GRAY)


def make_roadmap_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '演进路线与当前进度', size=32, bold=True, color=DARK_BG)

    phases = [
        ('Phase 0', '基础骨架', '✅', ACCENT_GREEN, '244 文件 / 38K 行\nGradle monorepo + 插件 + MCP + CLI'),
        ('Phase 1', 'Web IDE 实连', '✅', ACCENT_GREEN, '真流式 + Agentic Loop\n37 测试 + 跨栈画像'),
        ('Phase 1.5', '设计守护 + Docker', '✅', ACCENT_GREEN, '3 容器部署 + E2E 验证\n设计基线冻结'),
        ('Phase 1.6', 'AI 交付闭环 + SSO', '✅', ACCENT_GREEN, '9 工具 / 4 容器 / 130 测试\nKeycloak SSO + CRUD + 自动保存'),
        ('Phase 2', '单环运转', '👈', ACCENT_ORANGE, 'OODA 引擎 + SkillLoader 增强\n底线 CI + 用户试用'),
        ('Phase 3', '双环闭合', '⏳', TEXT_GRAY, 'ForgeNativeRuntime 完整\n进化环 + RuntimeAdapter'),
        ('Phase 4', '持续进化', '⏳', TEXT_GRAY, '多 Runtime + Skill 生态\n全组织部署'),
    ]

    # Timeline bar
    bar_y = Inches(1.8)
    add_shape(slide, Inches(0.8), bar_y + Inches(0.35), Inches(11.7), Inches(0.06), RGBColor(0xDD, 0xDD, 0xDD))

    for i, (name, keyword, status, color, desc) in enumerate(phases):
        x = Inches(0.6) + Inches(i * 1.75)

        # Dot on timeline
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.45), bar_y + Inches(0.18), Inches(0.4), Inches(0.4))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = dot.text_frame.paragraphs[0].add_run()
        run.text = status
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = 'Microsoft YaHei'

        # Phase name
        tb = add_textbox(slide, x, bar_y + Inches(0.7), Inches(1.6), Inches(0.5))
        set_text(tb.text_frame, name, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Keyword
        tb2 = add_textbox(slide, x, bar_y + Inches(1.1), Inches(1.6), Inches(0.4))
        set_text(tb2.text_frame, keyword, size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

        # Description card
        card_y = bar_y + Inches(1.7)
        card = add_shape(slide, x, card_y, Inches(1.6), Inches(2.7), LIGHT_GRAY)
        tb3 = add_textbox(slide, x + Inches(0.08), card_y + Inches(0.1), Inches(1.44), Inches(2.5))
        for line in desc.split('\n'):
            if tb3.text_frame.paragraphs[0].text == '':
                set_text(tb3.text_frame, line, size=11, color=TEXT_DARK)
            else:
                add_para(tb3.text_frame, line, size=11, color=TEXT_DARK, space_before=Pt(2), space_after=Pt(2))


def make_phase16_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, 'Phase 1.6 — AI 交付闭环 + SSO + UX 增强', size=32, bold=True, color=DARK_BG)
    add_para(tb.text_frame, '让 Web IDE 从"能对话"升级到"能交付"', size=18, color=TEXT_GRAY, space_before=Pt(4))

    features = [
        ('AI → Workspace\n交付闭环', 'workspace_write/read/list_files\nAI 直接读写 workspace 文件', ACCENT_BLUE),
        ('file_changed\n事件驱动', 'SSE 推送 → 文件树自动刷新\n编辑器自动打开新文件', ACCENT_BLUE),
        ('Keycloak SSO', 'OIDC PKCE 登录流程\n4 容器部署 + realm 导入', ACCENT_GREEN),
        ('Context Picker\n实连', '/api/context/search\n4 类别真实数据', ACCENT_GREEN),
        ('代码块 Apply', '一键写入 workspace\n20+ 语言扩展名映射', ACCENT_ORANGE),
        ('FileExplorer\nCRUD', '右键新建/重命名/删除\n文件 + 文件夹', ACCENT_ORANGE),
        ('未保存标记 +\n自动保存', 'Tab 蓝色圆点标记\n5 秒防抖自动保存', RGBColor(0x8B, 0x5C, 0xF6)),
        ('System Prompt\n交付指导', 'AI 被指导必须写文件\n而非仅在聊天展示代码', RGBColor(0x8B, 0x5C, 0xF6)),
    ]

    for i, (title, desc, color) in enumerate(features):
        col = i % 4
        row = i // 4
        x = Inches(0.8) + Inches(col * 3.1)
        y = Inches(1.8) + Inches(row * 2.6)

        # Card
        card = add_shape(slide, x, y, Inches(2.8), Inches(2.2), LIGHT_GRAY)
        # Color top bar
        add_shape(slide, x, y, Inches(2.8), Inches(0.06), color)

        tb = add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(1.9))
        set_text(tb.text_frame, title, size=16, bold=True, color=color)
        add_para(tb.text_frame, '', size=6)
        for line in desc.split('\n'):
            add_para(tb.text_frame, line, size=13, color=TEXT_GRAY, space_before=Pt(1), space_after=Pt(1))


def make_mcp_tools_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, 'MCP 工具体系 — 聚合架构', size=32, bold=True, color=DARK_BG)

    # Left: MCP Servers
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.3), RGBColor(0xE3, 0xF2, 0xFD))
    tb_l = add_textbox(slide, Inches(1), Inches(1.6), Inches(4.1), Inches(5))
    set_text(tb_l.text_frame, '后端 MCP Server（5 个 / 20 Tool）', size=18, bold=True, color=ACCENT_BLUE)
    servers = [
        ('forge-knowledge-mcp', '知识库读写 + 空白检测', '6 tools'),
        ('forge-database-mcp', '数据库元信息查询', '3 tools'),
        ('forge-service-graph-mcp', '服务依赖图谱', '5 tools'),
        ('forge-artifact-mcp', '制品库 + CVE 扫描', '3 tools'),
        ('forge-observability-mcp', '日志/指标/链路', '3 tools'),
    ]
    for name, desc, tools in servers:
        add_para(tb_l.text_frame, '', size=4)
        add_bullet(tb_l.text_frame, f'{name}', size=14, bold=True, color=TEXT_DARK)
        add_bullet(tb_l.text_frame, f'  {desc} ({tools})', size=12, color=TEXT_GRAY)

    # Arrow
    arrow_tb = add_textbox(slide, Inches(5.5), Inches(3.5), Inches(1.2), Inches(0.8))
    set_text(arrow_tb.text_frame, '→', size=36, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    desc_tb = add_textbox(slide, Inches(5.3), Inches(4.2), Inches(1.6), Inches(0.5))
    set_text(desc_tb.text_frame, 'McpProxyService\n聚合层', size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Right: Aggregated Tools
    add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xE8, 0xF5, 0xE9))
    tb_r = add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_r.text_frame, 'Claude API 暴露（9 个聚合 Tool）', size=18, bold=True, color=ACCENT_GREEN)

    tools = [
        ('knowledge_search', '统一知识搜索'),
        ('database_query', '数据库查询'),
        ('service_graph', '服务依赖查询'),
        ('artifact_search', '制品搜索'),
        ('observability_query', '可观测性查询'),
        ('codebase_analysis', '代码库分析'),
        ('workspace_write_file', 'AI 写文件到 workspace  ★'),
        ('workspace_read_file', 'AI 读 workspace 文件  ★'),
        ('workspace_list_files', 'AI 列出文件  ★'),
    ]
    for name, desc in tools:
        is_new = '★' in desc
        c = ACCENT_BLUE if is_new else TEXT_DARK
        add_bullet(tb_r.text_frame, f'{name}  —  {desc}', size=13, bold=is_new, color=c)

    # Design note
    note_tb = add_textbox(slide, Inches(7.2), Inches(6.2), Inches(5.4), Inches(0.5))
    set_text(note_tb.text_frame, '★ Phase 1.6 新增（AI 交付闭环核心能力）', size=12, bold=True, color=ACCENT_BLUE)


def make_tech_stack_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '技术架构与部署', size=32, bold=True, color=DARK_BG)

    # Left: Tech stack
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_GRAY)
    tb_l = add_textbox(slide, Inches(1), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_l.text_frame, 'Web IDE 技术选型', size=20, bold=True, color=ACCENT_BLUE)

    stack = [
        ('前端', 'React 19 + Next.js 15 (App Router)'),
        ('编辑器', 'Monaco Editor (VS Code 内核)'),
        ('UI', 'shadcn/ui + TailwindCSS'),
        ('状态', 'Zustand + React Query'),
        ('后端', 'Spring Boot 3 (Kotlin) + WebSocket'),
        ('认证', 'Keycloak 24.0 (OIDC PKCE)'),
        ('AI', 'Claude Agent SDK (流式 + MCP)'),
        ('Runtime', 'ForgeNativeRuntime (Phase 2+)'),
    ]
    for label, tech in stack:
        add_para(tb_l.text_frame, '', size=4)
        p = add_bullet(tb_l.text_frame, f'{label}', size=15, bold=True)
        run = p.add_run()
        run.text = f'  {tech}'
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT_GRAY
        run.font.name = 'Microsoft YaHei'

    # Right: Docker deployment
    add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xFF, 0xF8, 0xE1))
    tb_r = add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_r.text_frame, '4 容器部署架构', size=20, bold=True, color=ACCENT_ORANGE)

    containers = [
        ('Nginx', ':9000', '反向代理 + 统一入口'),
        ('Frontend', 'Next.js', 'SSR + 静态资源'),
        ('Backend', 'Spring Boot', 'API + SSE + MCP Proxy'),
        ('Keycloak', ':8180', 'SSO + OIDC + 用户管理'),
    ]
    for name, port, desc in containers:
        add_para(tb_r.text_frame, '', size=6)
        box_p = add_bullet(tb_r.text_frame, f'📦  {name}  ({port})', size=16, bold=True)
        add_bullet(tb_r.text_frame, f'     {desc}', size=14, color=TEXT_GRAY)

    add_para(tb_r.text_frame, '', size=12)
    add_para(tb_r.text_frame, '一键启动：', size=14, bold=True, color=TEXT_DARK)
    add_para(tb_r.text_frame, 'docker compose -f docker-compose.trial.yml up', size=13, color=ACCENT_BLUE)


def make_skill_system_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, 'Skill 体系 + 底线体系', size=32, bold=True, color=DARK_BG)

    # Left: Skill system
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xF3, 0xE5, 0xF5))
    tb_l = add_textbox(slide, Inches(1), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_l.text_frame, 'Skill 分层体系', size=20, bold=True, color=RGBColor(0x8B, 0x5C, 0xF6))

    skills = [
        ('交付阶段 Skill × 8', 'requirement-analysis / prd-writing / architecture-design\ndetailed-design / code-generation / test-case-writing\ntest-execution / deployment-ops'),
        ('Foundation Skill × 13', 'java-conventions / kotlin-conventions / spring-boot-patterns\ngradle-build / testing-standards / api-design\ndatabase-patterns / error-handling / logging-observability\nsecurity-practices / deployment-readiness-check\nenvironment-parity / design-baseline-guardian'),
        ('知识挖掘 Skill × 3', 'codebase-profiler（多语言）/ convention-miner（跨语言）'),
        ('跨栈迁移 Skill × 1', 'business-rule-extraction（业务规则 + 边界条件）'),
        ('Domain Skill × N', '按业务域扩展：payment / order / inventory ...'),
    ]

    for title, desc in skills:
        add_para(tb_l.text_frame, '', size=3)
        add_bullet(tb_l.text_frame, title, size=14, bold=True, color=TEXT_DARK)
        for line in desc.split('\n'):
            add_bullet(tb_l.text_frame, '  ' + line, size=11, color=TEXT_GRAY)

    # Right: Baseline system
    add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xFF, 0xEB, 0xEE))
    tb_r = add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_r.text_frame, '底线体系（8 个脚本）', size=20, bold=True, color=RED)

    baselines = [
        ('code-style-baseline', '阻断提交', 'ktlint + checkstyle 零违规'),
        ('security-baseline', '阻断提交', '无硬编码凭证 / SQL 注入 / XSS'),
        ('test-coverage-baseline', '阻断 merge', 'Service ≥ 80% 覆盖'),
        ('api-contract-baseline', '阻断 merge', 'OpenAPI Spec 一致'),
        ('architecture-baseline', '阻断 merge', '依赖方向正确'),
        ('deployment-preflight', '阻断部署', 'Lockfile + Dockerfile + 环境变量'),
        ('runtime-health', '阻断发布', '端点可达 + Bean 完整'),
        ('design-regression', '阻断 merge', 'API 契约 + UI 路由 + 数据模型'),
    ]

    for name, level, desc in baselines:
        add_para(tb_r.text_frame, '', size=2)
        p = add_bullet(tb_r.text_frame, f'{name}', size=13, bold=True, color=TEXT_DARK)
        add_bullet(tb_r.text_frame, f'  {level} — {desc}', size=11, color=TEXT_GRAY)

    add_para(tb_r.text_frame, '', size=8)
    add_para(tb_r.text_frame, '⚠ Gap：尚未集成到 CI', size=14, bold=True, color=ACCENT_ORANGE)


def make_gaps_next_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '已识别 Gap + Phase 2 规划', size=32, bold=True, color=DARK_BG)

    # Left: Gaps
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xFF, 0xF3, 0xE0))
    tb_l = add_textbox(slide, Inches(1), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_l.text_frame, '已识别 Gap（5 项）', size=20, bold=True, color=ACCENT_ORANGE)

    gaps = [
        ('底线 CI 集成', '8 个底线脚本已设计，尚未集成到 GitHub Actions', 'Phase 2'),
        ('Playwright E2E', '89 个手动验收用例未自动化', 'Phase 2'),
        ('内部用户试用', '3-5 人实际使用 ≥ 3 天尚未开展', 'Phase 2'),
        ('AI Chat 流式验证', '待 API Key 配置后验证完整流式体验', 'Phase 2'),
        ('进化环（环 2）', '全部组件待实现', 'Phase 3'),
    ]

    for title, desc, target in gaps:
        add_para(tb_l.text_frame, '', size=6)
        add_bullet(tb_l.text_frame, f'⚠  {title}', size=16, bold=True, color=ACCENT_ORANGE)
        add_bullet(tb_l.text_frame, f'    {desc}', size=13, color=TEXT_GRAY)
        add_bullet(tb_l.text_frame, f'    → 目标阶段: {target}', size=12, color=ACCENT_BLUE)

    # Right: Phase 2 plan
    add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.3), RGBColor(0xE3, 0xF2, 0xFD))
    tb_r = add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.4), Inches(5))
    set_text(tb_r.text_frame, 'Phase 2 关键交付物（12 项）', size=20, bold=True, color=ACCENT_BLUE)

    deliverables = [
        'SuperAgent OODA 引擎',
        'SkillLoader.kt 条件触发增强',
        'ProfileRouter.kt 后端智能路由',
        'MCP Server 完善（service-graph / database）',
        'convention-miner 跨语言增强',
        '13 个 Foundation Skill 全部实现',
        '端到端跨栈迁移 PoC（.NET → Java）',
        'agent-eval 评估集填充',
        '度量基线采集',
        '底线脚本 CI 集成',
        'Playwright E2E 自动化',
        '3-5 人内部试用 + 反馈',
    ]

    for d in deliverables:
        add_bullet(tb_r.text_frame, f'☐  {d}', size=13, color=TEXT_DARK)

    add_para(tb_r.text_frame, '', size=10)
    add_para(tb_r.text_frame, '验收：OODA 底线通过率 ≥ 70%', size=14, bold=True, color=ACCENT_GREEN)


def make_metrics_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_section_header(slide)

    tb = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '价值度量体系', size=32, bold=True, color=DARK_BG)

    metrics = [
        ('活跃度', '日活率 / Skill 触发 / MCP 调用 / Web IDE 会话', '4 周 60%+ 渗透', ACCENT_BLUE),
        ('效率', 'PR 周期 / 审查时间 / 新人上手 / Bug 修复', '缩短 30-60%', ACCENT_GREEN),
        ('质量', '生产 Bug 率 / 规范一致性 / 测试覆盖', '持续改善', ACCENT_GREEN),
        ('迁移效率', '跨栈迁移人天对比（传统 vs Forge）', '10x+ 人天提效', ACCENT_ORANGE),
        ('Agent 可靠性', 'OODA 一次通过率', '≥ 85%', ACCENT_BLUE),
        ('设计保真度', '设计基线回归测试通过率', '100%', RGBColor(0x8B, 0x5C, 0xF6)),
        ('部署效率', '首次部署成功率（pre-flight 后）', '≥ 90% (vs 7.7%)', RED),
        ('Runtime 独立性', 'Web IDE 不依赖 Claude Code 覆盖率', 'Phase 3 → 100%', TEXT_GRAY),
    ]

    for i, (name, indicator, target, color) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.8) + Inches(col * 3.1)
        y = Inches(1.5) + Inches(row * 2.8)

        card = add_shape(slide, x, y, Inches(2.8), Inches(2.4), LIGHT_GRAY)
        add_shape(slide, x, y, Inches(2.8), Inches(0.06), color)

        tb = add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(2.1))
        set_text(tb.text_frame, name, size=18, bold=True, color=color)
        add_para(tb.text_frame, '', size=4)
        add_para(tb.text_frame, indicator, size=12, color=TEXT_GRAY)
        add_para(tb.text_frame, '', size=8)
        add_para(tb.text_frame, f'目标: {target}', size=14, bold=True, color=TEXT_DARK)


def make_summary_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.12), ACCENT_BLUE)

    tb = add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1))
    set_text(tb.text_frame, '总结', size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = [
        ('✅  已完成', 'Phase 0 → 1 → 1.5 → 1.6，Web IDE 从骨架到可交付'),
        ('🔧  核心能力', '9 个 MCP 工具 / 5 个 Skill Profile / 13 Foundation Skill / 8 底线脚本'),
        ('📦  部署就绪', '4 容器 Docker 一键启动 + Keycloak SSO'),
        ('📊  质量保障', '130+ 测试 / 89 验收用例 / 设计基线 v5 冻结'),
        ('👈  下一步', 'Phase 2：OODA 引擎 + SkillLoader 增强 + 底线 CI + 用户试用'),
    ]

    for i, (title, desc) in enumerate(items):
        y = Inches(2.5) + Inches(i * 0.9)
        tb = add_textbox(slide, Inches(2), y, Inches(9), Inches(0.8))
        set_text(tb.text_frame, title, size=22, bold=True, color=ACCENT_GREEN)
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = f'    {desc}'
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        run.font.name = 'Microsoft YaHei'

    # Q&A
    tb = add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6))
    set_text(tb.text_frame, 'Q & A', size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Generate all slides ──────────────────────────────────────
make_title_slide()
make_goals_slide()
make_architecture_slide()
make_roadmap_slide()
make_phase16_slide()
make_mcp_tools_slide()
make_tech_stack_slide()
make_skill_system_slide()
make_gaps_next_slide()
make_metrics_slide()
make_summary_slide()

output_path = 'docs/forge-platform-baseline-v1.4.pptx'
prs.save(output_path)
print(f'✅ Generated: {output_path} ({len(prs.slides)} slides)')
