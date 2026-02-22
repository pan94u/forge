#!/usr/bin/env python3
"""
Forge Platform — Executive Presentation v2
Narrative: Problem → Insight → Vision → Proof → Value → Future
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Design System ─────────────────────────────────────────────
# Inspired by Apple keynote / Stripe deck aesthetics
BG_DARK = RGBColor(0x0F, 0x0F, 0x14)        # Almost black
BG_WARM = RGBColor(0x16, 0x16, 0x1D)        # Warm dark
BLUE = RGBColor(0x38, 0x7A, 0xFF)           # Electric blue
CYAN = RGBColor(0x00, 0xD4, 0xAA)           # Teal/cyan
ORANGE = RGBColor(0xFF, 0x9F, 0x43)         # Warm orange
PURPLE = RGBColor(0xA7, 0x7B, 0xFF)         # Soft purple
RED_SOFT = RGBColor(0xFF, 0x6B, 0x6B)       # Soft red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1 = RGBColor(0xE0, 0xE0, 0xE0)          # Light gray text
GRAY2 = RGBColor(0x8A, 0x8A, 0x9A)          # Medium gray
GRAY3 = RGBColor(0x50, 0x50, 0x60)          # Dark gray
CARD_BG = RGBColor(0x1E, 0x1E, 0x2A)        # Card background
CARD_BG2 = RGBColor(0x24, 0x24, 0x32)       # Lighter card
GREEN_BG = RGBColor(0x0A, 0x2A, 0x1F)       # Dark green tint
BLUE_BG = RGBColor(0x0A, 0x1A, 0x2F)        # Dark blue tint

# Fonts
FONT_CN = 'PingFang SC'
FONT_EN = 'SF Pro Display'
FONT_MONO = 'SF Mono'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ───────────────────────────────────────────────────
def dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK


def shape(slide, l, t, w, h, color, radius=None):
    if radius:
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        # Adjust corner radius via XML
        import lxml.etree as etree
        sp = s._element
        prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prstGeom is not None:
            avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            if avLst is None:
                avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
            gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    else:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def text(tf, txt, sz=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT_CN, spacing=None):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = Pt(spacing)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def para(tf, txt, sz=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT_CN, before=6, after=4, spacing=None):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if spacing:
        p.line_spacing = Pt(spacing)
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def accent_line(slide, x, y, w, color=BLUE):
    shape(slide, x, y, w, Inches(0.04), color)


def big_number(slide, x, y, number, label, color=CYAN):
    """Large metric display."""
    t = tb(slide, x, y, Inches(2.4), Inches(1.6))
    text(t.text_frame, number, sz=52, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_EN)
    para(t.text_frame, label, sz=14, color=GRAY2, align=PP_ALIGN.CENTER, before=2)


def card(slide, x, y, w, h, bg=CARD_BG):
    return shape(slide, x, y, w, h, bg, radius=8000)


# ── SLIDE 1: Title ────────────────────────────────────────────
def slide_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    # Subtle gradient accent at top
    shape(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

    # Main title
    t = tb(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5))
    text(t.text_frame, '重新定义软件交付', sz=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    t2 = tb(s, Inches(2), Inches(3.5), Inches(9), Inches(0.8))
    text(t2.text_frame, '当 1 个人 + AI 拥有 7 人团队的交付能力', sz=26, color=GRAY2, align=PP_ALIGN.CENTER)

    # Thin line
    accent_line(s, Inches(5.5), Inches(4.6), Inches(2.3), BLUE)

    # Project name
    t3 = tb(s, Inches(2), Inches(5.0), Inches(9), Inches(0.6))
    text(t3.text_frame, 'Forge 智能交付平台', sz=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Date
    t4 = tb(s, Inches(2), Inches(5.7), Inches(9), Inches(0.5))
    text(t4.text_frame, '2026.02', sz=16, color=GRAY3, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── SLIDE 2: The Problem ─────────────────────────────────────
def slide_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8))
    text(t.text_frame, '我们面对的现实', sz=36, bold=True, color=WHITE)

    problems = [
        ('5-7 人的 Scrum 团队', '每个角色一个人，知识在人脑里，人走知识就没了', RED_SOFT),
        ('AI 只是"更快的打字员"', '给每个人配一个 AI 助手 ≠ 交付效率的质变', ORANGE),
        ('跨栈迁移是黑洞', '没人懂旧技术栈，人工迁移动辄数十人月', ORANGE),
        ('经验不沉淀，错误在重复', '同样的 CORS 漏洞、同样的部署失败，每个项目重来一遍', RED_SOFT),
    ]

    for i, (title, desc, color) in enumerate(problems):
        y = Inches(1.8) + Inches(i * 1.3)

        # Number
        t_num = tb(s, Inches(1.2), y, Inches(0.6), Inches(0.6))
        text(t_num.text_frame, f'0{i+1}', sz=16, bold=True, color=GRAY3, font=FONT_EN)

        # Color bar
        shape(s, Inches(1.9), y + Inches(0.08), Inches(0.06), Inches(0.9), color)

        # Text
        t_title = tb(s, Inches(2.2), y, Inches(9), Inches(0.5))
        text(t_title.text_frame, title, sz=22, bold=True, color=WHITE)

        t_desc = tb(s, Inches(2.2), y + Inches(0.5), Inches(9), Inches(0.5))
        text(t_desc.text_frame, desc, sz=16, color=GRAY2)

    # Bottom punch
    t5 = tb(s, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6))
    text(t5.text_frame, '问题的本质不是工具不够，而是交付流程本身没有被重构。', sz=20, bold=True, color=CYAN)


# ── SLIDE 3: The Insight ──────────────────────────────────────
def slide_insight():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8))
    text(t.text_frame, '三个关键洞察', sz=36, bold=True, color=WHITE)

    insights = [
        (
            'SuperAgent，不是 Multi-Agent',
            '不是给 5 个角色各配一个 AI Agent。\n而是 1 个超级智能体，通过 Skill Profile 动态切换角色。\n一个人 + 一个 SuperAgent = 完整交付能力。',
            BLUE,
            '1 个智能体 > 5 个独立 Agent',
        ),
        (
            'Skill 是真正的护城河',
            '模型会被追平，Prompt 会被抄走。\n但编码了团队十年经验的 Skill 体系无法复制。\n这是碳基智能向硅基智能的系统性知识迁移。',
            CYAN,
            '25+ Skills = 团队经验的数字化资产',
        ),
        (
            '双环驱动：越用越好',
            '交付环解决"做什么"，进化环解决"越做越好"。\n每次交付产生数据 → 提取知识 → 反哺 AI → 下次更好。\n不是工具，是一个自我进化的交付系统。',
            PURPLE,
            '交付 × 进化 = 飞轮效应',
        ),
    ]

    for i, (title, desc, color, tagline) in enumerate(insights):
        x = Inches(0.8) + Inches(i * 4.1)
        y = Inches(1.8)

        c = card(s, x, y, Inches(3.8), Inches(4.8), CARD_BG)

        # Color top accent
        shape(s, x, y, Inches(3.8), Inches(0.06), color)

        # Content
        t_title = tb(s, x + Inches(0.3), y + Inches(0.4), Inches(3.2), Inches(0.8))
        text(t_title.text_frame, title, sz=20, bold=True, color=color)

        t_desc = tb(s, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(2.5))
        text(t_desc.text_frame, desc, sz=14, color=GRAY1, spacing=22)

        # Tagline
        accent_line(s, x + Inches(0.3), y + Inches(4.0), Inches(1.5), color)
        t_tag = tb(s, x + Inches(0.3), y + Inches(4.2), Inches(3.2), Inches(0.5))
        text(t_tag.text_frame, tagline, sz=13, bold=True, color=color)


# ── SLIDE 4: What SuperAgent Looks Like ───────────────────────
def slide_superagent():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8))
    text(t.text_frame, 'SuperAgent：一个智能体，五种角色', sz=36, bold=True, color=WHITE)

    t_sub = tb(s, Inches(1.2), Inches(1.2), Inches(10), Inches(0.5))
    text(t_sub.text_frame, '根据任务自动切换 Skill Profile，覆盖从需求到运维的完整交付链', sz=16, color=GRAY2)

    profiles = [
        ('规划', 'planning', '需求分析\nPRD 编写\n可行性评估', BLUE),
        ('设计', 'design', '架构设计\nAPI 设计\n影响分析', CYAN),
        ('开发', 'development', '代码生成\n规范遵循\n13 Foundation Skills', PURPLE),
        ('测试', 'testing', '用例设计\n边界识别\n自动化执行', ORANGE),
        ('运维', 'ops', '部署策略\n风险评估\nK8s 编排', RED_SOFT),
    ]

    for i, (name, profile, desc, color) in enumerate(profiles):
        x = Inches(0.6) + Inches(i * 2.52)
        y = Inches(2.2)

        c = card(s, x, y, Inches(2.3), Inches(3.5), CARD_BG)
        shape(s, x, y, Inches(2.3), Inches(0.06), color)

        # Profile name
        t_name = tb(s, x + Inches(0.2), y + Inches(0.3), Inches(1.9), Inches(0.5))
        text(t_name.text_frame, name, sz=24, bold=True, color=color, align=PP_ALIGN.CENTER)

        t_profile = tb(s, x + Inches(0.2), y + Inches(0.85), Inches(1.9), Inches(0.4))
        text(t_profile.text_frame, profile, sz=12, color=GRAY3, align=PP_ALIGN.CENTER, font=FONT_MONO)

        t_desc = tb(s, x + Inches(0.2), y + Inches(1.4), Inches(1.9), Inches(1.8))
        text(t_desc.text_frame, desc, sz=14, color=GRAY1, align=PP_ALIGN.CENTER, spacing=24)

    # OODA bar at bottom
    ooda_y = Inches(6.1)
    ooda_card = card(s, Inches(0.8), ooda_y, Inches(11.7), Inches(1.0), CARD_BG2)

    steps = [
        ('Observe', '观察上下文'),
        ('Orient', '分析 + 查知识库'),
        ('Decide', '制定方案'),
        ('Act', '执行 + 底线检查'),
        ('HITL', '人在回路审批'),
    ]
    for i, (en, cn) in enumerate(steps):
        sx = Inches(1.0) + Inches(i * 2.4)
        st = tb(s, sx, ooda_y + Inches(0.1), Inches(2), Inches(0.45))
        text(st.text_frame, en, sz=16, bold=True, color=CYAN, align=PP_ALIGN.CENTER, font=FONT_EN)
        st2 = tb(s, sx, ooda_y + Inches(0.5), Inches(2), Inches(0.4))
        text(st2.text_frame, cn, sz=12, color=GRAY2, align=PP_ALIGN.CENTER)

        if i < 4:
            arrow = tb(s, sx + Inches(1.85), ooda_y + Inches(0.15), Inches(0.5), Inches(0.4))
            text(arrow.text_frame, '→', sz=18, color=GRAY3, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── SLIDE 5: The Breakthrough ─────────────────────────────────
def slide_breakthrough():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    text(t.text_frame, '突破：AI 不再只是聊天，而是交付', sz=36, bold=True, color=WHITE)

    # Before / After comparison
    # BEFORE
    card(s, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), CARD_BG)
    shape(s, Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.06), RED_SOFT)

    t_b = tb(s, Inches(1.1), Inches(1.9), Inches(5), Inches(0.5))
    text(t_b.text_frame, 'Before：传统 AI 辅助', sz=20, bold=True, color=RED_SOFT)

    before_items = [
        '用户提问 → AI 在聊天窗口展示代码',
        '用户手动复制粘贴到文件',
        '用户手动创建/管理文件结构',
        '知识搜索返回 mock 数据',
        '没有认证，没有权限',
        'AI 不了解项目结构',
    ]
    t_items = tb(s, Inches(1.1), Inches(2.6), Inches(5), Inches(3.5))
    text(t_items.text_frame, '', sz=10)
    for item in before_items:
        para(t_items.text_frame, f'✗  {item}', sz=15, color=GRAY2, before=8, after=4)

    # Arrow
    a = tb(s, Inches(6.1), Inches(3.5), Inches(1.2), Inches(1))
    text(a.text_frame, '→', sz=48, bold=True, color=BLUE, align=PP_ALIGN.CENTER, font=FONT_EN)

    # AFTER
    card(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(5.2), CARD_BG)
    shape(s, Inches(6.9), Inches(1.6), Inches(5.9), Inches(0.06), CYAN)

    t_a = tb(s, Inches(7.2), Inches(1.9), Inches(5.3), Inches(0.5))
    text(t_a.text_frame, 'After：Forge AI 交付闭环', sz=20, bold=True, color=CYAN)

    after_items = [
        ('用户描述需求 → AI 直接写文件到 workspace', ''),
        ('文件树自动刷新，编辑器自动打开', 'SSE 事件驱动'),
        ('AI 主动了解项目结构再动手', 'workspace_list_files'),
        ('知识搜索返回真实文档', '12+ 知识库实连'),
        ('Keycloak SSO + OIDC PKCE', '企业级认证'),
        ('代码块一键 Apply，5 秒自动保存', '零摩擦体验'),
    ]
    t_items2 = tb(s, Inches(7.2), Inches(2.6), Inches(5.3), Inches(3.5))
    text(t_items2.text_frame, '', sz=10)
    for item, note in after_items:
        p = para(t_items2.text_frame, f'✓  {item}', sz=15, color=WHITE, before=8, after=2)
        if note:
            para(t_items2.text_frame, f'     {note}', sz=11, color=GRAY3, before=0, after=4)


# ── SLIDE 6: What We've Built (Numbers) ──────────────────────
def slide_numbers():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    text(t.text_frame, '我们已经走到了哪里', sz=36, bold=True, color=WHITE)

    t_sub = tb(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.5))
    text(t_sub.text_frame, 'Phase 0 → 1 → 1.5 → 1.6  |  从零到可交付的完整闭环', sz=16, color=GRAY2)

    # Big numbers row 1
    metrics1 = [
        ('45,000+', '行代码', CYAN),
        ('320+', '文件', BLUE),
        ('130+', '测试用例', PURPLE),
        ('89', '验收场景', ORANGE),
    ]

    for i, (num, label, color) in enumerate(metrics1):
        x = Inches(0.8) + Inches(i * 3.1)
        y = Inches(2.0)
        c = card(s, x, y, Inches(2.8), Inches(1.8), CARD_BG)

        t_num = tb(s, x, y + Inches(0.2), Inches(2.8), Inches(1.0))
        text(t_num.text_frame, num, sz=44, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = tb(s, x, y + Inches(1.15), Inches(2.8), Inches(0.4))
        text(t_label.text_frame, label, sz=14, color=GRAY2, align=PP_ALIGN.CENTER)

    # Big numbers row 2
    metrics2 = [
        ('9', 'MCP 工具', CYAN),
        ('5', 'Skill Profile', BLUE),
        ('13', 'Foundation Skill', PURPLE),
        ('4', 'Docker 容器', ORANGE),
    ]

    for i, (num, label, color) in enumerate(metrics2):
        x = Inches(0.8) + Inches(i * 3.1)
        y = Inches(4.1)
        c = card(s, x, y, Inches(2.8), Inches(1.8), CARD_BG)

        t_num = tb(s, x, y + Inches(0.2), Inches(2.8), Inches(1.0))
        text(t_num.text_frame, num, sz=44, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = tb(s, x, y + Inches(1.15), Inches(2.8), Inches(0.4))
        text(t_label.text_frame, label, sz=14, color=GRAY2, align=PP_ALIGN.CENTER)

    # Bottom tagline
    t_bottom = tb(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6))
    text(t_bottom.text_frame, '一键 docker compose up → 4 容器启动 → Keycloak SSO 登录 → AI 对话 → 代码直接写入项目', sz=15, color=GRAY3, align=PP_ALIGN.CENTER)


# ── SLIDE 7: Cross-stack Migration ────────────────────────────
def slide_migration():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    text(t.text_frame, '甜点场景：跨栈迁移 10x 提效', sz=36, bold=True, color=WHITE)

    t_sub = tb(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.5))
    text(t_sub.text_frame, '"没人懂旧技术栈"的困局，AI 可以系统性地解决', sz=16, color=GRAY2)

    # 5-step pipeline
    steps = [
        ('01', 'AI 代码考古', '解析 .NET/.sln/.csproj\n生成"项目说明书"', BLUE),
        ('02', '业务规则提取', '逐 Service 分析\n提取规则 + 边界条件', CYAN),
        ('03', '架构设计', '.NET 概念 → Java 映射\n依赖拓扑排列', PURPLE),
        ('04', '逐模块迁移', 'SuperAgent 全新编写\nFoundation Skills 保证规范', ORANGE),
        ('05', '行为等价验证', 'API 契约对比\n业务规则覆盖度检查', CYAN),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = Inches(0.5) + Inches(i * 2.52)
        y = Inches(2.2)

        c = card(s, x, y, Inches(2.3), Inches(3.0), CARD_BG)
        shape(s, x, y, Inches(2.3), Inches(0.06), color)

        t_num = tb(s, x + Inches(0.2), y + Inches(0.25), Inches(1.9), Inches(0.4))
        text(t_num.text_frame, num, sz=14, bold=True, color=GRAY3, font=FONT_EN)

        t_title = tb(s, x + Inches(0.2), y + Inches(0.65), Inches(1.9), Inches(0.5))
        text(t_title.text_frame, title, sz=18, bold=True, color=color)

        t_desc = tb(s, x + Inches(0.2), y + Inches(1.3), Inches(1.9), Inches(1.5))
        text(t_desc.text_frame, desc, sz=13, color=GRAY1, spacing=20)

        if i < 4:
            arrow = tb(s, x + Inches(2.15), y + Inches(1.0), Inches(0.6), Inches(0.5))
            text(arrow.text_frame, '→', sz=20, color=GRAY3, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Impact numbers
    impact_y = Inches(5.6)
    card(s, Inches(0.8), impact_y, Inches(11.7), Inches(1.3), GREEN_BG)

    impacts = [
        ('传统方式', '50-65 人天', RED_SOFT),
        ('→', '', GRAY3),
        ('Forge 辅助', '5-7 人天', CYAN),
        ('=', '', GRAY3),
        ('提效', '10-13x', ORANGE),
    ]
    for i, (label, value, color) in enumerate(impacts):
        ix = Inches(1.2) + Inches(i * 2.3)
        if label in ('→', '='):
            t_arr = tb(s, ix, impact_y + Inches(0.3), Inches(1), Inches(0.5))
            text(t_arr.text_frame, label, sz=28, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_EN)
        else:
            t_l = tb(s, ix, impact_y + Inches(0.15), Inches(2), Inches(0.45))
            text(t_l.text_frame, label, sz=13, color=GRAY2, align=PP_ALIGN.CENTER)
            t_v = tb(s, ix, impact_y + Inches(0.55), Inches(2), Inches(0.6))
            text(t_v.text_frame, value, sz=28, bold=True, color=color, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── SLIDE 8: Quality Safety Net ───────────────────────────────
def slide_quality():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    text(t.text_frame, '速度不缺，缺的是安全网', sz=36, bold=True, color=WHITE)

    t_sub = tb(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.6))
    text(t_sub.text_frame, '2400 行/小时的 AI 生产速度 + 8 道质量底线 = 又快又稳', sz=16, color=GRAY2)

    # The story
    card(s, Inches(0.8), Inches(2.0), Inches(6), Inches(2.2), CARD_BG)
    shape(s, Inches(0.8), Inches(2.0), Inches(6), Inches(0.06), RED_SOFT)
    t_story = tb(s, Inches(1.1), Inches(2.3), Inches(5.4), Inches(1.8))
    text(t_story.text_frame, '真实教训', sz=18, bold=True, color=RED_SOFT)
    para(t_story.text_frame, 'Phase 1.5 首次 Docker 部署：13 次尝试，90 分钟，10 个计划外问题。', sz=15, color=GRAY1, before=10)
    para(t_story.text_frame, 'CORS 通配符直接进入 commit，无人察觉。', sz=15, color=GRAY1, before=4)
    para(t_story.text_frame, '→ 首次部署成功率仅 7.7%（1/13）', sz=16, bold=True, color=RED_SOFT, before=10)

    # The solution
    card(s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(2.2), CARD_BG)
    shape(s, Inches(7.2), Inches(2.0), Inches(5.6), Inches(0.06), CYAN)
    t_sol = tb(s, Inches(7.5), Inches(2.3), Inches(5), Inches(1.8))
    text(t_sol.text_frame, '系统性解决', sz=18, bold=True, color=CYAN)
    para(t_sol.text_frame, '痛点编码为 Skill + 底线脚本，让未来的开发者无感获得保护。', sz=15, color=GRAY1, before=10)
    para(t_sol.text_frame, '不依赖模型"自觉"，而是自动化质量关卡。', sz=15, color=GRAY1, before=4)
    para(t_sol.text_frame, '→ 目标首次部署成功率 ≥ 90%', sz=16, bold=True, color=CYAN, before=10)

    # 8 baselines visualization
    baselines = [
        ('代码规范', '阻断提交'),
        ('安全扫描', '阻断提交'),
        ('测试覆盖', '阻断 merge'),
        ('API 契约', '阻断 merge'),
        ('架构守护', '阻断 merge'),
        ('部署预检', '阻断部署'),
        ('运行健康', '阻断发布'),
        ('设计回归', '阻断 merge'),
    ]

    for i, (name, level) in enumerate(baselines):
        x = Inches(0.8) + Inches(i * 1.56)
        y = Inches(4.6)
        c = card(s, x, y, Inches(1.4), Inches(1.5), CARD_BG)
        shape(s, x, y, Inches(1.4), Inches(0.04), CYAN if i < 5 else ORANGE)

        t_name = tb(s, x + Inches(0.08), y + Inches(0.2), Inches(1.24), Inches(0.6))
        text(t_name.text_frame, name, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        t_level = tb(s, x + Inches(0.08), y + Inches(0.85), Inches(1.24), Inches(0.4))
        text(t_level.text_frame, level, sz=11, color=GRAY3, align=PP_ALIGN.CENTER)

    # Principle
    t_p = tb(s, Inches(0.8), Inches(6.4), Inches(12), Inches(0.6))
    text(t_p.text_frame, '原则：已验证的设计不可隐式退化  |  开发者痛点即平台能力', sz=15, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ── SLIDE 9: Strategic Roadmap ────────────────────────────────
def slide_roadmap():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    text(t.text_frame, '从工具到平台到生态', sz=36, bold=True, color=WHITE)

    # Timeline
    phases = [
        ('已完成', [
            ('Phase 0-1', '骨架 + Web IDE 实连', '✓'),
            ('Phase 1.5', '设计守护 + Docker', '✓'),
            ('Phase 1.6', 'AI 交付闭环 + SSO', '✓'),
        ], CYAN),
        ('下一步', [
            ('Phase 2', '单环运转\nOODA 引擎 + SkillLoader\n底线 CI + 用户试用', '→'),
        ], ORANGE),
        ('远景', [
            ('Phase 3', '双环闭合\n进化环 + Runtime 抽象\nClaude Code 变为可选', '○'),
            ('Phase 4', '持续进化\n多 Runtime + Skill 生态\n全组织部署', '○'),
        ], GRAY3),
    ]

    y_base = Inches(1.6)

    # Section labels and content
    x_cursor = Inches(0.8)
    for section_name, items, color in phases:
        # Section header
        t_sec = tb(s, x_cursor, y_base, Inches(len(items) * 3.2), Inches(0.5))
        text(t_sec.text_frame, section_name, sz=14, bold=True, color=color, font=FONT_EN)
        accent_line(s, x_cursor, y_base + Inches(0.45), Inches(len(items) * 3.0), color)

        for i, (name, desc, icon) in enumerate(items):
            x = x_cursor + Inches(i * 3.2)
            y = y_base + Inches(0.7)

            c = card(s, x, y, Inches(2.9), Inches(4.2), CARD_BG)
            shape(s, x, y, Inches(2.9), Inches(0.05), color)

            # Icon
            t_icon = tb(s, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.4))
            text(t_icon.text_frame, icon, sz=18, bold=True, color=color, font=FONT_EN)

            # Name
            t_name = tb(s, x + Inches(0.7), y + Inches(0.2), Inches(2), Inches(0.5))
            text(t_name.text_frame, name, sz=18, bold=True, color=WHITE)

            # Desc
            t_desc = tb(s, x + Inches(0.2), y + Inches(0.85), Inches(2.5), Inches(3.0))
            text(t_desc.text_frame, desc, sz=13, color=GRAY1, spacing=22)

        x_cursor += Inches(len(items) * 3.2) + Inches(0.3)

    # Vision statement
    t_v = tb(s, Inches(1), Inches(6.3), Inches(11), Inches(0.8))
    text(t_v.text_frame, '终局愿景：1 人 + SuperAgent 全自主交付  |  Skill 生态开放  |  知识持续进化', sz=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ── SLIDE 10: The Ask ─────────────────────────────────────────
def slide_ask():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    t = tb(s, Inches(1.2), Inches(0.6), Inches(10), Inches(0.8))
    text(t.text_frame, '下一步需要什么', sz=36, bold=True, color=WHITE)

    asks = [
        (
            'API Key 配置 + 试用环境',
            '完成 Claude API Key 接入，搭建内部试用环境，3-5 人真实使用 ≥ 3 天。\n当前 AI 交付闭环已就绪，差这一步就能跑通完整流程。',
            CYAN,
            '紧迫度：高',
        ),
        (
            '1 个试点项目（跨栈迁移优先）',
            '选择一个小型 .NET → Java 迁移模块，验证 10x 提效假设。\n这是 Forge 最有说服力的场景，可产出量化 ROI 数据。',
            ORANGE,
            '价值：最高',
        ),
        (
            'CI/CD 管道接入',
            '将 8 个底线脚本集成到 GitHub Actions，让质量安全网自动生效。\n从"设计了"到"跑起来"，需要 DevOps 协助。',
            PURPLE,
            '依赖：DevOps',
        ),
    ]

    for i, (title, desc, color, tag) in enumerate(asks):
        y = Inches(1.8) + Inches(i * 1.8)

        c = card(s, Inches(0.8), y, Inches(11.7), Inches(1.5), CARD_BG)
        shape(s, Inches(0.8), y, Inches(0.08), Inches(1.5), color)

        # Number
        t_num = tb(s, Inches(1.1), y + Inches(0.2), Inches(0.8), Inches(0.5))
        text(t_num.text_frame, f'0{i+1}', sz=28, bold=True, color=color, font=FONT_EN)

        # Title
        t_title = tb(s, Inches(2), y + Inches(0.15), Inches(7), Inches(0.5))
        text(t_title.text_frame, title, sz=22, bold=True, color=WHITE)

        # Desc
        t_desc = tb(s, Inches(2), y + Inches(0.65), Inches(7.5), Inches(0.7))
        text(t_desc.text_frame, desc, sz=14, color=GRAY2, spacing=20)

        # Tag
        tag_box = card(s, Inches(10.5), y + Inches(0.3), Inches(1.8), Inches(0.45), color)
        tag_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tag_box.text_frame.paragraphs[0].add_run()
        r.text = tag
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE if color != ORANGE else BG_DARK
        r.font.name = FONT_CN

    # Summary
    t_sum = tb(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6))
    text(t_sum.text_frame, 'Phase 2 目标：OODA 底线通过率 ≥ 70%  |  跨栈迁移 PoC 业务规则覆盖 ≥ 90%  |  3+ 人完成试用', sz=14, color=GRAY3, align=PP_ALIGN.CENTER)


# ── SLIDE 11: Closing ─────────────────────────────────────────
def slide_closing():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    dark_bg(s)

    shape(s, Inches(0), Inches(0), W, Inches(0.05), BLUE)

    # Core message
    t = tb(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2))
    text(t.text_frame, '不是给每个人配一个 AI 助手', sz=32, color=GRAY3, align=PP_ALIGN.CENTER)

    t2 = tb(s, Inches(1.5), Inches(3.2), Inches(10), Inches(1.2))
    text(t2.text_frame, '而是用 AI 重构交付流程本身', sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    accent_line(s, Inches(5.5), Inches(4.6), Inches(2.3), CYAN)

    t3 = tb(s, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6))
    text(t3.text_frame, 'Forge — 让 1 个人拥有 7 人团队的交付能力', sz=20, color=CYAN, align=PP_ALIGN.CENTER)

    t4 = tb(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6))
    text(t4.text_frame, 'Thank You', sz=24, bold=True, color=GRAY3, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── Generate ──────────────────────────────────────────────────
slide_title()        # 1. 封面
slide_problem()      # 2. 问题
slide_insight()      # 3. 洞察
slide_superagent()   # 4. SuperAgent
slide_breakthrough() # 5. 突破（Before/After）
slide_numbers()      # 6. 数据
slide_migration()    # 7. 跨栈迁移
slide_quality()      # 8. 质量安全网
slide_roadmap()      # 9. 路线图
slide_ask()          # 10. 下一步
slide_closing()      # 11. 结语

output = 'docs/forge-platform-executive-v2.pptx'
prs.save(output)
print(f'✅ Generated: {output} ({len(prs.slides)} slides)')
