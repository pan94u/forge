#!/usr/bin/env python3
"""
Forge Platform — Apple Keynote Style Executive Presentation
Design: Pure black, white typography, single accent, extreme restraint
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import lxml.etree as etree

# ── Apple-inspired palette (极度克制) ─────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TITLE = RGBColor(0xF5, 0xF5, 0xF7)    # Apple headline gray-white
GRAY_BODY = RGBColor(0x86, 0x86, 0x8B)      # Apple body gray
GRAY_DIM = RGBColor(0x48, 0x48, 0x4A)       # Dimmed text
GRAY_CARD = RGBColor(0x1D, 0x1D, 0x1F)      # Apple dark card
ACCENT = RGBColor(0x29, 0x97, 0xFF)         # Apple blue — 唯一强调色，少用

FONT = 'PingFang SC'
FONT_EN = 'SF Pro Display'
FONT_MONO = 'SF Mono'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def bg_black(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = BLACK


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def rounded(slide, l, t, w, h, color, r=6000):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    sp = s._element
    prstGeom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if prstGeom is not None:
        avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        if avLst is None:
            avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
        else:
            for child in list(avLst):
                avLst.remove(child)
        gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
        gd.set('name', 'adj')
        gd.set('fmla', f'val {r}')
    return s


def txbox(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def txt(tf, s, sz=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT, lh=None):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if lh:
        p.line_spacing = Pt(lh)
    r = p.add_run()
    r.text = s
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def add(tf, s, sz=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font=FONT, before=6, after=2, lh=None):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if lh:
        p.line_spacing = Pt(lh)
    r = p.add_run()
    r.text = s
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def draw_logo(slide, cx, cy, scale=1.0):
    """
    Forge logo: abstract anvil + spark.
    An anvil silhouette (trapezoid base + rectangular body)
    with a small diamond spark above — symbolizing craftsmanship + AI spark.
    """
    s = scale
    # Anvil body — wide base
    base = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 42*s), int(cy + 10*s),
        int(84*s), int(18*s)
    )
    base.fill.solid()
    base.fill.fore_color.rgb = GRAY_TITLE
    base.line.fill.background()

    # Anvil neck — narrow middle
    neck = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 18*s), int(cy - 14*s),
        int(36*s), int(28*s)
    )
    neck.fill.solid()
    neck.fill.fore_color.rgb = GRAY_TITLE
    neck.line.fill.background()

    # Anvil top — working surface (wider than neck)
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx - 34*s), int(cy - 24*s),
        int(68*s), int(14*s)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = GRAY_TITLE
    top.line.fill.background()

    # Anvil horn — right extension
    horn = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(cx + 34*s), int(cy - 20*s),
        int(22*s), int(8*s)
    )
    horn.fill.solid()
    horn.fill.fore_color.rgb = GRAY_TITLE
    horn.line.fill.background()

    # Spark — small diamond above the anvil
    spark = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        int(cx - 7*s), int(cy - 44*s),
        int(14*s), int(14*s)
    )
    spark.fill.solid()
    spark.fill.fore_color.rgb = ACCENT
    spark.line.fill.background()

    # Small spark particles
    for dx, dy, sz_s in [(-18, -36, 5), (20, -38, 4), (-8, -52, 3), (14, -50, 3)]:
        p = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(cx + dx*s), int(cy + dy*s),
            int(sz_s*s), int(sz_s*s)
        )
        p.fill.solid()
        p.fill.fore_color.rgb = ACCENT
        p.line.fill.background()


def draw_logo_small(slide, cx, cy, scale=0.5):
    """Simplified small logo mark — just spark diamond."""
    spark = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        int(cx - 10*scale*2), int(cy - 10*scale*2),
        int(20*scale*2), int(20*scale*2)
    )
    spark.fill.solid()
    spark.fill.fore_color.rgb = ACCENT
    spark.line.fill.background()


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Title — minimal, centered, logo
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s01_title():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    # Logo
    draw_logo(s, Inches(6.666), Inches(2.2), scale=Emu(Inches(0.012)))

    # FORGE wordmark
    t = txbox(s, Inches(2), Inches(3.2), Inches(9.3), Inches(1.0))
    txt(t.text_frame, 'FORGE', sz=52, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Tagline
    t2 = txbox(s, Inches(2), Inches(4.2), Inches(9.3), Inches(0.6))
    txt(t2.text_frame, '面向 AI 时代的智能交付平台', sz=22, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    # Thin line
    rect(s, Inches(5.8), Inches(5.1), Inches(1.7), Inches(0.02), GRAY_DIM)

    # Date
    t3 = txbox(s, Inches(2), Inches(5.4), Inches(9.3), Inches(0.5))
    txt(t3.text_frame, '2026.02', sz=15, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: One line hook
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s02_hook():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5))
    txt(t.text_frame, '当 1 个人 + AI\n拥有 7 人团队的交付能力', sz=52, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, lh=68)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: The Problem — 4 pain points, clean layout
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s03_problem():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '我们面对的现实', sz=40, bold=True, color=GRAY_TITLE)

    problems = [
        ('5-7 人的 Scrum 团队', '每个角色一个人，知识在人脑里，人走知识就没了'),
        ('AI 只是"更快的打字员"', '给每个人配一个 AI 助手，不等于交付效率的质变'),
        ('跨栈迁移是黑洞', '没人懂旧技术栈，人工迁移动辄数十人月'),
        ('经验不沉淀，错误在重复', '同样的部署失败、同样的安全漏洞，每个项目重来一遍'),
    ]

    for i, (title, desc) in enumerate(problems):
        y = Inches(2.0) + Inches(i * 1.25)

        # Thin left accent
        rect(s, Inches(1.2), y + Inches(0.05), Inches(0.04), Inches(0.85), GRAY_DIM)

        t_title = txbox(s, Inches(1.6), y, Inches(10), Inches(0.45))
        txt(t_title.text_frame, title, sz=22, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, Inches(1.6), y + Inches(0.5), Inches(10), Inches(0.4))
        txt(t_desc.text_frame, desc, sz=16, color=GRAY_BODY)

    # Bottom punch — the ONLY use of accent color on this slide
    t5 = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.5))
    txt(t5.text_frame, '问题的本质不是工具不够，是交付流程本身没有被重构。', sz=18, color=ACCENT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: Three Insights — 3 cards, restrained color
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s04_insight():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '三个关键洞察', sz=40, bold=True, color=GRAY_TITLE)

    insights = [
        (
            'SuperAgent\n不是 Multi-Agent',
            '不是 5 个角色各配一个 AI。\n而是 1 个超级智能体，\n通过 Skill Profile 动态切换。',
            '1 个智能体 > 5 个独立 Agent',
        ),
        (
            'Skill 是\n真正的护城河',
            '模型会被追平，Prompt 会被抄走。\n但编码了十年经验的 Skill 体系，\n无法复制。',
            '25+ Skills = 经验的数字化资产',
        ),
        (
            '双环驱动\n越用越好',
            '交付环做事，进化环学习。\n每次交付产生知识沉淀，\n反哺下一次交付。',
            '交付 × 进化 = 飞轮效应',
        ),
    ]

    for i, (title, desc, tagline) in enumerate(insights):
        x = Inches(0.8) + Inches(i * 4.15)
        y = Inches(1.9)

        c = rounded(s, x, y, Inches(3.85), Inches(4.6), GRAY_CARD)

        # Thin top line — only accent element
        rect(s, x, y, Inches(3.85), Inches(0.03), ACCENT if i == 0 else GRAY_DIM)

        t_title = txbox(s, x + Inches(0.4), y + Inches(0.5), Inches(3.05), Inches(1.2))
        txt(t_title.text_frame, title, sz=24, bold=True, color=GRAY_TITLE, lh=32)

        t_desc = txbox(s, x + Inches(0.4), y + Inches(1.8), Inches(3.05), Inches(1.8))
        txt(t_desc.text_frame, desc, sz=15, color=GRAY_BODY, lh=24)

        # Tagline at bottom
        t_tag = txbox(s, x + Inches(0.4), y + Inches(3.9), Inches(3.05), Inches(0.4))
        txt(t_tag.text_frame, tagline, sz=13, color=GRAY_DIM)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: SuperAgent — 5 profiles + OODA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s05_superagent():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '一个智能体，五种角色', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.3), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '根据任务自动切换 Skill Profile，覆盖从需求到运维的完整交付链', sz=16, color=GRAY_BODY)

    profiles = [
        ('规划', 'Planning', '需求分析\nPRD 编写'),
        ('设计', 'Design', '架构设计\nAPI 设计'),
        ('开发', 'Development', '代码生成\n13 Foundation Skills'),
        ('测试', 'Testing', '用例设计\n自动化执行'),
        ('运维', 'Ops', '部署策略\nK8s 编排'),
    ]

    for i, (cn, en, desc) in enumerate(profiles):
        x = Inches(0.6) + Inches(i * 2.52)
        y = Inches(2.2)

        c = rounded(s, x, y, Inches(2.3), Inches(3.0), GRAY_CARD)
        # Only first card gets accent line
        rect(s, x, y, Inches(2.3), Inches(0.03), ACCENT if i == 2 else GRAY_DIM)

        t_cn = txbox(s, x + Inches(0.25), y + Inches(0.35), Inches(1.8), Inches(0.5))
        txt(t_cn.text_frame, cn, sz=26, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER)

        t_en = txbox(s, x + Inches(0.25), y + Inches(0.85), Inches(1.8), Inches(0.35))
        txt(t_en.text_frame, en, sz=12, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)

        t_desc = txbox(s, x + Inches(0.25), y + Inches(1.5), Inches(1.8), Inches(1.2))
        txt(t_desc.text_frame, desc, sz=14, color=GRAY_BODY, align=PP_ALIGN.CENTER, lh=22)

    # OODA bar
    ooda_y = Inches(5.7)
    rounded(s, Inches(0.8), ooda_y, Inches(11.7), Inches(1.2), GRAY_CARD)

    t_ooda_label = txbox(s, Inches(1.2), ooda_y + Inches(0.1), Inches(3), Inches(0.35))
    txt(t_ooda_label.text_frame, 'OODA 循环', sz=12, color=GRAY_DIM, font=FONT_EN)

    steps = ['Observe  观察', 'Orient  分析', 'Decide  决策', 'Act  执行', 'HITL  人在回路']
    for i, step in enumerate(steps):
        sx = Inches(1.0) + Inches(i * 2.35)
        st = txbox(s, sx, ooda_y + Inches(0.45), Inches(2.1), Inches(0.5))
        txt(st.text_frame, step, sz=15, color=GRAY_TITLE if i < 4 else ACCENT, align=PP_ALIGN.CENTER)

        if i < 4:
            arrow = txbox(s, sx + Inches(1.9), ooda_y + Inches(0.45), Inches(0.5), Inches(0.5))
            txt(arrow.text_frame, '→', sz=16, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: Four-Layer Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s06_architecture():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    txt(t.text_frame, '四层解耦，独立演进', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '每一层可以独立升级，不影响其他层 — 这是平台能活过模型迭代的根本原因', sz=16, color=GRAY_BODY)

    layers = [
        (
            '用户交互层',
            'User Interaction',
            'CLI / VS Code (开发者) + Web IDE (全角色)',
            'Monaco 编辑器  ·  AI Chat  ·  知识浏览  ·  可视化工作流',
            False,
        ),
        (
            'SuperAgent 层',
            'Intelligence',
            '5 Profiles × 25+ Skills × 8 Baselines',
            'Skill Profile Router  ·  Baseline Runner  ·  Skill Market ↔ Capability Market (MCP)',
            True,   # accent — this is the core
        ),
        (
            'Agent Runtime 层',
            'Runtime',
            'MCP 协议  ·  Tool Calling  ·  沙箱  ·  可观测',
            '9 聚合工具  ·  安全围栏  ·  SuperAgent 评估框架',
            False,
        ),
        (
            '数据与模型层',
            'Data & Model',
            '公域模型 (Claude)  +  私域数据 (知识库 12+)  +  垂域模型',
            'MCP Gateway → Knowledge / Database / ServiceGraph / Artifact / Observability',
            False,
        ),
    ]

    for i, (cn, en, line1, line2, is_accent) in enumerate(layers):
        y = Inches(1.75) + Inches(i * 1.35)

        # Card
        c = rounded(s, Inches(0.8), y, Inches(11.7), Inches(1.2), GRAY_CARD)

        # Left accent strip
        rect(s, Inches(0.8), y, Inches(0.05), Inches(1.2), ACCENT if is_accent else GRAY_DIM)

        # Layer number circle — just text
        t_num = txbox(s, Inches(1.15), y + Inches(0.15), Inches(0.65), Inches(0.55))
        txt(t_num.text_frame, f'L{i+1}', sz=14, bold=True,
            color=ACCENT if is_accent else GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

        # Chinese name
        t_cn = txbox(s, Inches(1.9), y + Inches(0.1), Inches(2.5), Inches(0.4))
        txt(t_cn.text_frame, cn, sz=20, bold=True, color=GRAY_TITLE)

        # English label
        t_en = txbox(s, Inches(4.3), y + Inches(0.15), Inches(2), Inches(0.35))
        txt(t_en.text_frame, en, sz=11, color=GRAY_DIM, font=FONT_MONO)

        # Description line 1
        t_l1 = txbox(s, Inches(1.9), y + Inches(0.55), Inches(10), Inches(0.3))
        txt(t_l1.text_frame, line1, sz=14, color=ACCENT if is_accent else GRAY_BODY)

        # Description line 2
        t_l2 = txbox(s, Inches(1.9), y + Inches(0.85), Inches(10), Inches(0.3))
        txt(t_l2.text_frame, line2, sz=12, color=GRAY_DIM)

        # Arrow between layers
        if i < 3:
            arr_y = y + Inches(1.2)
            t_arr = txbox(s, Inches(6.2), arr_y, Inches(0.9), Inches(0.2))
            txt(t_arr.text_frame, '↓', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Bottom insight
    t_bot = txbox(s, Inches(1.2), Inches(7.0), Inches(11), Inches(0.4))
    txt(t_bot.text_frame, '模型从 Opus 换成 GPT?  换 L3 的 Adapter。Skills 和 Baselines 一行不改。', sz=14, color=ACCENT, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Stable vs Volatile — Adapter Isolation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s07_stable_volatile():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.5), Inches(10), Inches(0.8))
    txt(t.text_frame, '稳态固化 · 敏态抽象', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.1), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '模型一定会被追平，能力编码为 Skill 才是真正的护城河', sz=16, color=GRAY_BODY)

    # ── Left: 稳态 ──
    rounded(s, Inches(0.8), Inches(1.8), Inches(4.8), Inches(4.8), GRAY_CARD)
    rect(s, Inches(0.8), Inches(1.8), Inches(4.8), Inches(0.03), ACCENT)

    t_sl = txbox(s, Inches(1.2), Inches(2.1), Inches(4), Inches(0.5))
    txt(t_sl.text_frame, '稳态 — 核心资产', sz=20, bold=True, color=ACCENT)

    t_sld = txbox(s, Inches(1.2), Inches(2.6), Inches(4), Inches(0.3))
    txt(t_sld.text_frame, '越积越厚，无法复制', sz=13, color=GRAY_DIM)

    stable_items = [
        ('25+ Skills', '十年工程经验的数字化编码'),
        ('8 Baselines', '质量底线自动守护'),
        ('5 Profiles', '全交付角色覆盖'),
        ('12+ Knowledge', '私域知识持续沉淀'),
        ('OODA 流程', '标准化决策环路'),
    ]

    for i, (item, desc) in enumerate(stable_items):
        iy = Inches(3.1) + Inches(i * 0.64)
        t_item = txbox(s, Inches(1.5), iy, Inches(2.2), Inches(0.35))
        txt(t_item.text_frame, item, sz=17, bold=True, color=GRAY_TITLE, font=FONT_EN)
        t_desc = txbox(s, Inches(3.6), iy + Inches(0.02), Inches(2), Inches(0.35))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_BODY)

    # ── Center: Adapter ──
    mid_x = Inches(5.95)
    mid_y = Inches(3.2)

    # Adapter box
    adapter = rounded(s, mid_x - Inches(0.55), mid_y, Inches(1.1), Inches(2.0), GRAY_DIM, r=4000)

    t_ad = txbox(s, mid_x - Inches(0.55), mid_y + Inches(0.2), Inches(1.1), Inches(1.6))
    tf = t_ad.text_frame
    tf.word_wrap = True
    txt(tf, 'A', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN)
    add(tf, 'd', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'a', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'p', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 't', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'e', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)
    add(tf, 'r', sz=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font=FONT_EN, before=2, after=2)

    # Arrows left → adapter
    t_al = txbox(s, Inches(5.6), Inches(3.9), Inches(0.5), Inches(0.4))
    txt(t_al.text_frame, '→', sz=18, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Arrows adapter → right
    t_ar = txbox(s, Inches(6.9), Inches(3.9), Inches(0.5), Inches(0.4))
    txt(t_ar.text_frame, '→', sz=18, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # ── Right: 敏态 ──
    rounded(s, Inches(7.7), Inches(1.8), Inches(4.8), Inches(4.8), GRAY_CARD)
    rect(s, Inches(7.7), Inches(1.8), Inches(4.8), Inches(0.03), GRAY_DIM)

    t_vl = txbox(s, Inches(8.1), Inches(2.1), Inches(4), Inches(0.5))
    txt(t_vl.text_frame, '敏态 — 可替换组件', sz=20, bold=True, color=GRAY_BODY)

    t_vld = txbox(s, Inches(8.1), Inches(2.6), Inches(4), Inches(0.3))
    txt(t_vld.text_frame, '技术迭代时只换这里', sz=13, color=GRAY_DIM)

    volatile_items = [
        ('ModelAdapter', 'Claude → GPT → Local'),
        ('RuntimeAdapter', 'Claude Code → ForgeNative'),
        ('AssetFormatAdapter', 'SKILL.md v1 → v2 → ...'),
    ]

    for i, (item, desc) in enumerate(volatile_items):
        iy = Inches(3.1) + Inches(i * 0.8)
        t_item = txbox(s, Inches(8.3), iy, Inches(4), Inches(0.35))
        txt(t_item.text_frame, item, sz=16, bold=True, color=GRAY_TITLE, font=FONT_MONO)
        t_desc = txbox(s, Inches(8.3), iy + Inches(0.35), Inches(4), Inches(0.3))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_DIM)

    # Adapter labels at bottom
    adapters_label = [
        'ModelAdapter.kt  ·  RuntimeAdapter.kt  ·  AssetFormatAdapter.kt'
    ]
    t_alab = txbox(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.3))
    txt(t_alab.text_frame, adapters_label[0], sz=12, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_MONO)

    # Bottom insight — the key message for leadership
    rounded(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.85), GRAY_CARD)
    rect(s, Inches(0.8), Inches(6.2), Inches(0.05), Inches(0.85), ACCENT)

    t_bot = txbox(s, Inches(1.3), Inches(6.35), Inches(10.8), Inches(0.6))
    txt(t_bot.text_frame, '战略意义', sz=13, bold=True, color=ACCENT)
    add(t_bot.text_frame,
        'OpenAI 发布更好的模型？换一个 Adapter。Claude Code 被替代？切换 ForgeNativeRuntime。\n'
        '但 25 个 Skill 编码的十年工程经验，竞争对手无法复制。',
        sz=14, color=GRAY_BODY, before=4, lh=22)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: The Breakthrough — Before / After (was 6)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s08_breakthrough():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, 'AI 不再只是聊天，而是交付', sz=40, bold=True, color=GRAY_TITLE)

    # BEFORE
    rounded(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), GRAY_CARD)
    rect(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.03), GRAY_DIM)

    t_b = txbox(s, Inches(1.2), Inches(2.1), Inches(4.8), Inches(0.5))
    txt(t_b.text_frame, 'Before', sz=18, bold=True, color=GRAY_DIM, font=FONT_EN)

    before = [
        'AI 在聊天窗口展示代码',
        '用户手动复制粘贴到文件',
        '用户手动创建和管理文件',
        '知识搜索返回 mock 数据',
        '没有认证，没有权限控制',
        'AI 不了解项目结构',
    ]
    t_items = txbox(s, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.5))
    txt(t_items.text_frame, '', sz=8)
    for item in before:
        add(t_items.text_frame, f'  {item}', sz=16, color=GRAY_BODY, before=10, after=2)

    # Arrow
    a = txbox(s, Inches(6.2), Inches(3.6), Inches(0.9), Inches(0.8))
    txt(a.text_frame, '→', sz=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    # AFTER
    rounded(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(4.8), GRAY_CARD)
    rect(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(0.03), ACCENT)

    t_a = txbox(s, Inches(7.3), Inches(2.1), Inches(5.1), Inches(0.5))
    txt(t_a.text_frame, 'After — Forge', sz=18, bold=True, color=ACCENT, font=FONT_EN)

    after = [
        ('AI 直接写文件到 workspace', 'workspace_write / read / list'),
        ('文件树自动刷新，编辑器自动打开', 'SSE file_changed 事件驱动'),
        ('AI 主动了解项目结构再动手', 'workspace_list_files'),
        ('知识搜索返回真实文档', '12+ 知识库文档实连'),
        ('Keycloak SSO 企业级认证', 'OIDC PKCE + 4 容器部署'),
        ('代码块一键 Apply + 5 秒自动保存', '零摩擦交付体验'),
    ]
    t_items2 = txbox(s, Inches(7.3), Inches(2.7), Inches(5.1), Inches(3.5))
    txt(t_items2.text_frame, '', sz=8)
    for item, note in after:
        add(t_items2.text_frame, f'  {item}', sz=16, color=GRAY_TITLE, before=6, after=0)
        add(t_items2.text_frame, f'     {note}', sz=11, color=GRAY_DIM, before=0, after=4, font=FONT_MONO)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Numbers — big, clean, Apple-style
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s09_numbers():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '我们走到了哪里', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.25), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, 'Phase 0 → 1 → 1.5 → 1.6  完成', sz=16, color=GRAY_BODY)

    metrics = [
        ('45,000+', '行代码'),
        ('320+', '文件'),
        ('130+', '测试'),
        ('89', '验收场景'),
    ]

    for i, (num, label) in enumerate(metrics):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(2.4)

        t_num = txbox(s, x, y, Inches(2.9), Inches(1.5))
        txt(t_num.text_frame, num, sz=64, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = txbox(s, x, y + Inches(1.3), Inches(2.9), Inches(0.4))
        txt(t_label.text_frame, label, sz=16, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    metrics2 = [
        ('9', 'MCP 工具'),
        ('5', 'Skill Profile'),
        ('13', 'Foundation Skill'),
        ('4', 'Docker 容器'),
    ]

    for i, (num, label) in enumerate(metrics2):
        x = Inches(0.8) + Inches(i * 3.15)
        y = Inches(4.5)

        t_num = txbox(s, x, y, Inches(2.9), Inches(1.3))
        txt(t_num.text_frame, num, sz=56, bold=True, color=ACCENT if i == 0 else GRAY_TITLE, align=PP_ALIGN.CENTER, font=FONT_EN)

        t_label = txbox(s, x, y + Inches(1.1), Inches(2.9), Inches(0.4))
        txt(t_label.text_frame, label, sz=16, color=GRAY_BODY, align=PP_ALIGN.CENTER)

    # Thin separator
    rect(s, Inches(1.2), Inches(4.3), Inches(10.9), Inches(0.01), GRAY_DIM)

    # Bottom
    t_bot = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4))
    txt(t_bot.text_frame, 'docker compose up → Keycloak SSO → AI 对话 → 代码写入 workspace → 文件树刷新', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: Cross-stack Migration — the sweet spot
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s10_migration():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '跨栈迁移：10 倍提效', sz=40, bold=True, color=GRAY_TITLE)

    t_sub = txbox(s, Inches(1.2), Inches(1.3), Inches(10), Inches(0.4))
    txt(t_sub.text_frame, '"没人懂旧技术栈"的困局，AI 可以系统性地解决', sz=16, color=GRAY_BODY)

    steps = [
        ('01', 'AI 代码考古', '解析项目结构\n生成项目说明书'),
        ('02', '业务规则提取', '逐模块分析\n规则 + 边界条件'),
        ('03', '架构设计', '.NET → Java 映射\n依赖拓扑排列'),
        ('04', '逐模块迁移', 'SuperAgent 编写\n规范自动遵循'),
        ('05', '行为等价验证', 'API 契约对比\n规则覆盖度检查'),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = Inches(0.5) + Inches(i * 2.52)
        y = Inches(2.2)

        c = rounded(s, x, y, Inches(2.3), Inches(2.6), GRAY_CARD)
        rect(s, x, y, Inches(2.3), Inches(0.03), ACCENT if i == 0 else GRAY_DIM)

        t_num = txbox(s, x + Inches(0.2), y + Inches(0.25), Inches(0.5), Inches(0.35))
        txt(t_num.text_frame, num, sz=13, color=GRAY_DIM, font=FONT_EN)

        t_title = txbox(s, x + Inches(0.2), y + Inches(0.65), Inches(1.9), Inches(0.5))
        txt(t_title.text_frame, title, sz=17, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, x + Inches(0.2), y + Inches(1.2), Inches(1.9), Inches(1.2))
        txt(t_desc.text_frame, desc, sz=13, color=GRAY_BODY, lh=20)

        if i < 4:
            arrow = txbox(s, x + Inches(2.15), y + Inches(0.8), Inches(0.55), Inches(0.5))
            txt(arrow.text_frame, '→', sz=18, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Impact comparison — big numbers
    y2 = Inches(5.2)
    rect(s, Inches(1.2), y2, Inches(10.9), Inches(0.01), GRAY_DIM)

    # Traditional
    t_trad_l = txbox(s, Inches(1.5), y2 + Inches(0.3), Inches(3), Inches(0.3))
    txt(t_trad_l.text_frame, '传统方式', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)
    t_trad_n = txbox(s, Inches(1.5), y2 + Inches(0.6), Inches(3), Inches(0.8))
    txt(t_trad_n.text_frame, '50-65 人天', sz=32, bold=True, color=GRAY_BODY, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Arrow
    t_arr = txbox(s, Inches(5), y2 + Inches(0.6), Inches(1), Inches(0.6))
    txt(t_arr.text_frame, '→', sz=28, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Forge
    t_forge_l = txbox(s, Inches(6.2), y2 + Inches(0.3), Inches(3), Inches(0.3))
    txt(t_forge_l.text_frame, 'Forge 辅助', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)
    t_forge_n = txbox(s, Inches(6.2), y2 + Inches(0.6), Inches(3), Inches(0.8))
    txt(t_forge_n.text_frame, '5-7 人天', sz=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)

    # Multiplier
    t_mult = txbox(s, Inches(9.8), y2 + Inches(0.4), Inches(2.5), Inches(1.0))
    txt(t_mult.text_frame, '10x', sz=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_EN)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9: Quality — the safety net story
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s11_quality():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '速度不缺，缺的是安全网', sz=40, bold=True, color=GRAY_TITLE)

    # The real story
    t_story = txbox(s, Inches(1.2), Inches(1.8), Inches(10), Inches(1.8))
    txt(t_story.text_frame, '首次部署成功率', sz=14, color=GRAY_DIM)
    add(t_story.text_frame, '7.7%', sz=72, bold=True, color=GRAY_TITLE, font=FONT_EN, before=0, after=0)
    add(t_story.text_frame, '13 次尝试，90 分钟，10 个计划外问题。CORS 通配符直接进入 commit，无人察觉。', sz=16, color=GRAY_BODY, before=8)

    # Thin line
    rect(s, Inches(1.2), Inches(4.0), Inches(10.9), Inches(0.01), GRAY_DIM)

    # Solution
    t_sol = txbox(s, Inches(1.2), Inches(4.3), Inches(5.5), Inches(1.8))
    txt(t_sol.text_frame, '系统性解决', sz=20, bold=True, color=GRAY_TITLE)
    add(t_sol.text_frame, '痛点不修补，而是编码为平台能力。', sz=16, color=GRAY_BODY, before=10)
    add(t_sol.text_frame, '8 道质量底线自动拦截，', sz=16, color=GRAY_BODY, before=4)
    add(t_sol.text_frame, '让未来的开发者无感地获得保护。', sz=16, color=GRAY_BODY, before=4)

    # Target
    t_target = txbox(s, Inches(7.5), Inches(4.3), Inches(5), Inches(1.8))
    txt(t_target.text_frame, '目标成功率', sz=14, color=GRAY_DIM)
    add(t_target.text_frame, '≥ 90%', sz=56, bold=True, color=ACCENT, font=FONT_EN, before=0)

    # Principle
    t_p = txbox(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4))
    txt(t_p.text_frame, '已验证的设计不可隐式退化  ·  开发者痛点即平台能力', sz=14, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10: Roadmap — past / next / vision
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s12_roadmap():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '从工具到平台到生态', sz=40, bold=True, color=GRAY_TITLE)

    # Three columns: Done / Next / Vision
    cols = [
        ('已完成', ACCENT, [
            ('Phase 0', '骨架搭建'),
            ('Phase 1', 'Web IDE 实连'),
            ('Phase 1.5', '设计守护 + Docker'),
            ('Phase 1.6', 'AI 交付闭环 + SSO'),
        ]),
        ('下一步', GRAY_TITLE, [
            ('Phase 2', 'OODA 引擎 + SkillLoader'),
            ('', '底线脚本 CI 集成'),
            ('', '跨栈迁移 PoC'),
            ('', '3-5 人内部试用'),
        ]),
        ('远景', GRAY_DIM, [
            ('Phase 3', '进化环闭合'),
            ('', 'Runtime 完全抽象'),
            ('Phase 4', 'Skill 生态开放'),
            ('', '全组织部署'),
        ]),
    ]

    for ci, (section, color, items) in enumerate(cols):
        x = Inches(0.8) + Inches(ci * 4.15)
        y = Inches(1.8)

        # Section label
        t_sec = txbox(s, x, y, Inches(3.85), Inches(0.5))
        txt(t_sec.text_frame, section, sz=14, bold=True, color=color, font=FONT_EN)

        rect(s, x, y + Inches(0.5), Inches(3.85), Inches(0.02), color)

        c = rounded(s, x, y + Inches(0.7), Inches(3.85), Inches(4.3), GRAY_CARD)

        t_items = txbox(s, x + Inches(0.35), y + Inches(1.0), Inches(3.15), Inches(3.8))
        txt(t_items.text_frame, '', sz=8)

        for phase, desc in items:
            if phase:
                add(t_items.text_frame, phase, sz=15, bold=True, color=GRAY_TITLE, before=14, after=2)
                add(t_items.text_frame, desc, sz=14, color=GRAY_BODY, before=2, after=2)
            else:
                add(t_items.text_frame, desc, sz=14, color=GRAY_BODY, before=6, after=2)

    # Vision statement
    t_v = txbox(s, Inches(1.2), Inches(6.6), Inches(11), Inches(0.4))
    txt(t_v.text_frame, '终局：1 人 + SuperAgent 全自主交付  ·  Skill 生态开放  ·  知识持续进化', sz=15, color=GRAY_DIM, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11: The Ask — what we need
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s13_ask():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    t = txbox(s, Inches(1.2), Inches(0.7), Inches(10), Inches(0.8))
    txt(t.text_frame, '下一步需要什么', sz=40, bold=True, color=GRAY_TITLE)

    asks = [
        ('01', 'API Key + 试用环境', '完成 Claude API 接入，3-5 人真实使用。\n交付闭环已就绪，差这一步跑通完整流程。'),
        ('02', '1 个试点项目', '选一个小型 .NET → Java 模块，验证 10x 假设。\n最有说服力的场景，可产出量化 ROI。'),
        ('03', 'CI/CD 管道接入', '8 个底线脚本集成到 GitHub Actions。\n让质量安全网从"设计了"变成"跑起来"。'),
    ]

    for i, (num, title, desc) in enumerate(asks):
        y = Inches(1.9) + Inches(i * 1.7)

        rounded(s, Inches(0.8), y, Inches(11.7), Inches(1.4), GRAY_CARD)
        # Left accent bar
        rect(s, Inches(0.8), y, Inches(0.05), Inches(1.4), ACCENT if i == 0 else GRAY_DIM)

        t_num = txbox(s, Inches(1.2), y + Inches(0.25), Inches(0.8), Inches(0.5))
        txt(t_num.text_frame, num, sz=28, bold=True, color=ACCENT if i == 0 else GRAY_DIM, font=FONT_EN)

        t_title = txbox(s, Inches(2.1), y + Inches(0.15), Inches(5), Inches(0.5))
        txt(t_title.text_frame, title, sz=24, bold=True, color=GRAY_TITLE)

        t_desc = txbox(s, Inches(2.1), y + Inches(0.65), Inches(9), Inches(0.6))
        txt(t_desc.text_frame, desc, sz=14, color=GRAY_BODY, lh=22)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12: Closing — one powerful sentence
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def s14_closing():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_black(s)

    # Logo small
    draw_logo(s, Inches(6.666), Inches(1.8), scale=Emu(Inches(0.008)))

    # Dimmed line
    t1 = txbox(s, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.8))
    txt(t1.text_frame, '不是给每个人配一个 AI 助手', sz=28, color=GRAY_DIM, align=PP_ALIGN.CENTER)

    # Bold statement
    t2 = txbox(s, Inches(1.5), Inches(3.8), Inches(10.3), Inches(1.0))
    txt(t2.text_frame, '而是用 AI 重构交付流程本身', sz=42, bold=True, color=GRAY_TITLE, align=PP_ALIGN.CENTER)

    # Subtle line
    rect(s, Inches(5.8), Inches(5.1), Inches(1.7), Inches(0.02), GRAY_DIM)

    # Thank you
    t3 = txbox(s, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5))
    txt(t3.text_frame, 'Thank You', sz=20, color=GRAY_DIM, align=PP_ALIGN.CENTER, font=FONT_EN)


# ── Generate ──────────────────────────────────────────────────
s01_title()
s02_hook()
s03_problem()
s04_insight()
s05_superagent()
s06_architecture()
s07_stable_volatile()
s08_breakthrough()
s09_numbers()
s10_migration()
s11_quality()
s12_roadmap()
s13_ask()
s14_closing()

out = 'docs/forge-platform-executive-v3.pptx'
prs.save(out)
print(f'✅ {out} ({len(prs.slides)} slides)')
